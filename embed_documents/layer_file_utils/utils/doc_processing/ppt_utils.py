from pptx import Presentation
from typing import Union
import io

def extract_text_from_ppt(file_input: Union[str, io.BytesIO]) -> str:
    """
    Extract text from PowerPoint presentation while preserving structure.
    Handles both file paths and file-like objects.
    
    Args:
        file_input (Union[str, io.BytesIO]): Path to the PPT file or file-like object
    
    Returns:
        str: Extracted text with preserved structure
    """
    try:
        # Read PowerPoint presentation
        if isinstance(file_input, str):
            prs = Presentation(file_input)
        else:
            prs = Presentation(file_input)
        
        text = ""
        for slide_number, slide in enumerate(prs.slides, 1):
            text += f"\nSlide {slide_number}:\n"
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
            
            # Extract text from tables
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        text += row_text + "\n"
                    text += "\n"
        
        return text
                    
    except Exception as e:
        raise Exception(f"Error extracting text from PowerPoint: {str(e)}") 